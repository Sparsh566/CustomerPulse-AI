import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_slide(prs, layout_idx, title_text, content_text=None, content2_text=None):
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    if slide.shapes.title:
        slide.shapes.title.text = title_text
        for par in slide.shapes.title.text_frame.paragraphs:
            par.font.bold = True
            par.font.color.rgb = RGBColor(0, 51, 102)

    # Content 1
    if content_text and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = content_text
        for par in tf.paragraphs:
            par.font.size = Pt(16)
            
    # Content 2
    if content2_text and len(slide.placeholders) > 2:
        tf2 = slide.placeholders[2].text_frame
        tf2.text = content2_text
        for par in tf2.paragraphs:
            par.font.size = Pt(16)
            
    return slide

prs = Presentation()

# --- Slide 1: Title ---
slide_layout = prs.slide_layouts[0] # Title slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "CustomerPulse AI"
subtitle.text = "Unified AI-Powered Complaint Intelligence Platform for BFSI\n\nIntake • Triage • Emotion Timeline • Urgency Score • SLA Risk • Resolution • Customer 360 • Reporting"
for par in title.text_frame.paragraphs:
    par.font.bold = True
    par.font.color.rgb = RGBColor(0, 51, 102)

# --- Slide 2: The Problem ---
content = (
    "The Complaint Chaos in Banking\n"
    "• Complaints arrive via 6+ channels with zero unified view\n"
    "• Manual triage leads to wrong priority, wrong team, SLA missed\n"
    "• Agents work completely blind — no customer history before responding\n"
    "• Genuine emergencies treated same as routine requests\n"
    "• No emotion tracking — agents can't see if customer is about to leave\n"
    "• No duplicate detection — same issue resolved multiple times, wasting effort\n"
    "• RBI regulatory reporting done manually — error-prone and time-consuming\n"
    "• Repeat complaints from same customer go unnoticed — no pattern detection\n\n"
    "Key Statistics:\n"
    "• 68% of customers leave after a poorly handled complaint\n"
    "• 40% of agent time is wasted on manual triage\n"
    "• 3x higher churn risk when SLAs are missed\n"
    "• 1 in 4 CRITICAL complaints gets zero special treatment\n"
)
add_slide(prs, 1, "The Problem", content)

# --- Slide 3: Why Existing Tools Fail ---
content = (
    "Why Existing Platforms Fall Short:\n"
    "• Salesforce Service Cloud: Expensive, no RBI reporting, no emotion AI, no predictive SLA\n"
    "• Freshdesk: No urgency detection, no customer 360, no Gen-AI draft responses\n"
    "• Zoho Desk: No sentiment drift, no compliance layer, no duplicate clustering\n\n"
    "What CustomerPulse AI Solves:\n"
    "✓ BFSI-Native — built for RBI / IRDAI compliance from day one\n"
    "✓ Customer Emotion Timeline — tracks mood per message\n"
    "✓ Financial Urgency Score — genuine priorities for emergencies\n"
    "✓ Customer 360 — full history profile for the agent\n"
    "✓ Agent Whisper — real-time AI coaching during live interactions\n"
)
add_slide(prs, 1, "Why Existing Tools Fail BFSI", content)

# --- Slide 4: Solution Overview ---
content = (
    "One Platform. Full Complaint Lifecycle.\n\n"
    "01. Complaint Intake\n"
    "    Omnichannel ingestion, Auto ticket ID\n\n"
    "02. AI Engine\n"
    "    NLP-based sentiment, category, auto-routing\n\n"
    "03. Emotion & Urgency Intelligence [NEW]\n"
    "    Emotion Timeline, Financial Urgency, Agent Whisper\n\n"
    "04. Customer 360 Profile [NEW]\n"
    "    Full complaint history, Churn risk score\n\n"
    "05. Agent Action Layer\n"
    "    AI-drafted response, One-click escalation\n\n"
    "06. Reporting & Compliance\n"
    "    RBI Annexure/SLA reports, automated exports\n"
)
add_slide(prs, 1, "Solution Overview", content)

# --- Slide 5: System Architecture ---
content = (
    "How It All Works Together\n\n"
    "1. INPUT CHANNELS\n"
    "   Email, Live Chat, Branch, Social Media, WhatsApp (Proposed)\n\n"
    "2. COMPLAINT INTAKE & VALIDATION\n"
    "   Auto Ticket ID, Zod Validation, Channel Tag\n\n"
    "3. AI ENGINE (Gemini Flash via OpenRouter)\n"
    "   Sentiment Analysis, Categorization, Auto-Routing\n\n"
    "4. AI INTELLIGENCE LAYER [Core Differentiator]\n"
    "   Emotion Timeline, Financial Urgency Score, Agent Whisper\n\n"
    "5. DATA BACKEND & INFRASTRUCTURE\n"
    "   Supabase Postgres, Edge Functions, Audit Logging\n\n"
    "6. OUTPUT LAYER\n"
    "   Role-Based Dashboard, RBI Reports, Customer SLA Monitor\n"
)
add_slide(prs, 1, "System Architecture", content)

# --- Slide 6: Core AI & NLP Engine ---
content = (
    "Analytical AI Capabilities:\n"
    "• Auto-Categorization: Product & complaint classification with manual override\n"
    "• Duplicate & Cluster Detection: Semantic similarity to prevent redundant effort\n"
    "• Sentiment & Severity Scoring: Tracking sentiment drift and rating priority 1-10\n"
    "• Predictive SLA Breach Detection: Auto-reassign if breach risk > 80%\n\n"
    "Generative AI Capabilities:\n"
    "• AI Draft Responses: Under 60 seconds customized by tone\n"
    "• Next-Best Actions & Similar Cases: Fast recommendations based on history\n"
    "• Regulatory Report Generation: One-click formatted data exports\n"
    "• AI Analytics Assistant: Query datasets with natural language\n"
)
add_slide(prs, 1, "Core AI & NLP Features", content)

# --- Slide 7: Our Edge: 3 New Features ---
content = (
    "Feature 1: Customer Emotion Timeline\n"
    "• Tracks emotional journey (Calm → Anxious → Frustrated)\n"
    "• Includes Empathy Coach & De-escalation Trigger\n"
    "Feature 2: Financial Urgency Score Engine\n"
    "• Scores 0-100 on real financial/medical harm\n"
    "• 8 Genuine Urgency Categories (Critical/High/Medium)\n"
    "Feature 3: Customer 360 Profile\n"
    "• Complete history visible to agent before responding\n"
    "• Repeat pattern alerts & predictive churn risk\n"
    "Bonus: Agent Whisper\n"
    "• Real-time, invisible suggestions to agent during chat\n"
)
add_slide(prs, 1, "What Makes CustomerPulse Different", content)

# --- Slide 8: Dashboard & User Experience ---
content = (
    "Role-Based Dashboards for Agent, Manager, Admin & Customer:\n\n"
    "1. Operations Dashboard\n"
    "   KPI cards, predictive SLA surface, and Burst alerts\n\n"
    "2. Complaint Detail View\n"
    "   Emotion timeline, Whisper panel, Unified thread\n\n"
    "3. Customer 360 Profile View\n"
    "   Complete history summary & agent notes\n\n"
    "4. Analytics & Reports\n"
    "   RBI/SLA reports, Urgency metrics, AI natural language query\n\n"
    "5. Public Tracking & Performance\n"
    "   Customer-facing tracking; Agent empathy matching score\n"
)
add_slide(prs, 1, "Dashboard & User Experience", content)

# --- Slide 9: Impact & Business Case ---
content = (
    "Key Operational Metrics (Before vs After):\n"
    "• Avg Resolution Time: ~18 hours (63% faster than 48-72h avg)\n"
    "• SLA Compliance: ~92% (up from 60%)\n"
    "• Agent Triage Time: Reduced by 75%\n"
    "• Critical Urgency Response: < 2 hours guaranteed\n\n"
    "Business Case & Potential:\n"
    "• Target Market: 100K+ bank branches & 2K+ NBFCs in India\n"
    "• Revenue Model: SaaS Subscriptions + Premium compliance reporting\n"
    "• The 'Why Now': RBI mandates are strict, complaint volume is growing\n"
)
add_slide(prs, 1, "Impact & Business Case", content)

# --- Slide 10: Tech Stack + Roadmap ---
content = (
    "Modern Tech Stack:\n"
    "• Frontend: React 18, Vite, Tailwind CSS, shadcn/ui\n"
    "• Backend: Supabase Postgres & Edge Functions (Deno)\n"
    "• AI Layer: Gemini Flash via OpenRouter + Custom Edge Pipelines\n\n"
    "Roadmap:\n"
    "• Phase 1: WhatsApp API, Multilingual NLP, Auto-escalation\n"
    "• Phase 2: Predictive churn score, Core banking integration\n\n"
    "Summary:\n"
    "CustomerPulse AI — turning every complaint into an insight, "
    "every emotion into an action, every emergency into a priority, "
    "and every SLA into a promise kept.\n"
)
add_slide(prs, 1, "Built to Scale — Stack & Roadmap", content)

prs.save('CustomerPulse_AI_Pitch.pptx')
print("Presentation generated successfully: CustomerPulse_AI_Pitch.pptx")
